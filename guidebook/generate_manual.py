"""Generate QC Suite Operator Manual Book (DOCX + PPTX) — Blue-White Theme."""
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
import os

GUIDEBOOK_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "guidebook")

# ── Color Palette ──
BLUE_DARK = RGBColor(0x1A, 0x36, 0x5C)
BLUE_MID = RGBColor(0x2B, 0x6C, 0xB0)
BLUE_LIGHT = RGBColor(0x63, 0xB3, 0xED)
BLUE_BG = RGBColor(0xEB, 0xF4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x20, 0x2C)
GRAY = RGBColor(0x71, 0x80, 0x96)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT_GREEN = RGBColor(0x27, 0x67, 0x49)
ACCENT_RED = RGBColor(0x9B, 0x2C, 0x2C)
ACCENT_YELLOW = RGBColor(0xB7, 0x79, 0x1E)

P_BLUE_DARK = PptxRGB(0x1A, 0x36, 0x5C)
P_BLUE_MID = PptxRGB(0x2B, 0x6C, 0xB0)
P_BLUE_LIGHT = PptxRGB(0x63, 0xB3, 0xED)
P_BLUE_BG = PptxRGB(0xEB, 0xF4, 0xFF)
P_WHITE = PptxRGB(0xFF, 0xFF, 0xFF)
P_DARK = PptxRGB(0x1A, 0x20, 0x2C)
P_GREEN = PptxRGB(0x27, 0x67, 0x49)
P_RED = PptxRGB(0x9B, 0x2C, 0x2C)
P_GRAY = PptxRGB(0x71, 0x80, 0x96)
P_LIGHT_GRAY = PptxRGB(0xE2, 0xE8, 0xF0)

_ROI_ALLOWED = {"x", "y", "w", "h", "rotation", "width"}


def _color_to_hex(color):
    if isinstance(color, str):
        return color
    try:
        return f"{color.red:02X}{color.green:02X}{color.blue:02X}"
    except AttributeError:
        return "FFFFFF"


def set_cell_shading(cell, color):
    hex_color = _color_to_hex(color)
    shading = cell._element.get_or_add_tcPr()
    shading_elem = shading.makeelement(qn('w:shd'), {
        qn('w:val'): 'clear', qn('w:color'): 'auto', qn('w:fill'): hex_color,
    })
    shading.append(shading_elem)


def add_heading(doc, text, level=1, color=BLUE_DARK):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return h


def add_para(doc, text, bold=False, size=10, color=DARK_TEXT, align=None):
    p = doc.add_paragraph()
    if align:
        p.alignment = align
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return p


def add_bullet(doc, text, size=10, color=DARK_TEXT):
    p = doc.add_paragraph(style='List Bullet')
    run = p.add_run(text)
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return p


def add_table(doc, headers, rows, header_color=BLUE_MID, header_text=WHITE):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    header_hex = _color_to_hex(header_color)
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ''
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(h)
        run.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = header_text
        run.font.name = 'Calibri'
        set_cell_shading(cell, header_hex)
    alt_hex = _color_to_hex(LIGHT_GRAY)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.rows[ri + 1].cells[ci]
            cell.text = ''
            p = cell.paragraphs[0]
            run = p.add_run(str(val))
            run.font.size = Pt(9)
            run.font.color.rgb = DARK_TEXT
            run.font.name = 'Calibri'
            if ri % 2 == 1:
                set_cell_shading(cell, alt_hex)
    return table


def add_photo_placeholder(doc, caption, width=Inches(5.5), height=Inches(3.0)):
    """Add a photo placeholder box with caption."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    # Draw a bordered box using a table with 1 cell
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.rows[0].cells[0]
    cell.width = width
    # Set cell borders
    tc = cell._element
    tcPr = tc.get_or_add_tcPr()
    tcBorders = tcPr.makeelement(qn('w:tcBorders'), {})
    for border_name in ['top', 'left', 'bottom', 'right']:
        border = tcBorders.makeelement(qn(f'w:{border_name}'), {
            qn('w:val'): 'single', qn('w:sz'): '6',
            qn('w:space'): '0', qn('w:color'): 'B0C4DE',
        })
        tcBorders.append(border)
    tcPr.append(tcBorders)
    # Fill light blue
    set_cell_shading(cell, 'EBF4FF')
    # Placeholder text
    p_cell = cell.paragraphs[0]
    p_cell.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_cell.add_run('[ FOTO DI SINI ]')
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
    run.font.name = 'Calibri'
    # Set cell height
    tr = cell._element.getparent()
    trPr = tr.get_or_add_trPr()
    trHeight = trPr.makeelement(qn('w:trHeight'), {
        qn('w:val'): str(int(height / 0.014 * 720)),  # approximate twips
        qn('w:hRule'): 'exact',
    })
    trPr.append(trHeight)
    # Caption below
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(caption)
    run.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY
    run.font.name = 'Calibri'
    return table


def add_page_break(doc):
    doc.add_page_break()


def add_section_header(doc, number, title):
    """Add a styled section header with blue bar."""
    p = doc.add_paragraph()
    run = p.add_run(f"BAB {number}")
    run.font.size = Pt(14)
    run.bold = True
    run.font.color.rgb = WHITE
    run.font.name = 'Calibri'
    # Add blue line below
    p2 = doc.add_paragraph()
    run2 = p2.add_run("━" * 50)
    run2.font.color.rgb = BLUE_MID
    run2.font.size = Pt(8)
    add_heading(doc, title, level=1)


# ═══════════════════════════════════════════════════════════════════
# DOCX GENERATION
# ═══════════════════════════════════════════════════════════════════

def generate_docx():
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(10.5)

    # ── Cover Page ──
    for _ in range(4):
        doc.add_paragraph()

    # Logo placeholder
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("[ LOGO ASKI ]")
    run.font.size = Pt(16)
    run.font.color.rgb = GRAY
    run.font.name = 'Calibri'

    doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("QC SUITE")
    run.font.size = Pt(40)
    run.bold = True
    run.font.color.rgb = BLUE_DARK
    run.font.name = 'Calibri'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Operator Manual Book")
    run.font.size = Pt(22)
    run.font.color.rgb = BLUE_MID
    run.font.name = 'Calibri'

    doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Sistem Inspeksi QC Sticker & Component Counter")
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY
    run.font.name = 'Calibri'

    for _ in range(4):
        doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("PT ASKI")
    run.font.size = Pt(16)
    run.bold = True
    run.font.color.rgb = BLUE_DARK
    run.font.name = 'Calibri'

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Versi 1.0 — 2026")
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY
    run.font.name = 'Calibri'

    add_page_break(doc)

    # ── Table of Contents ──
    add_heading(doc, "Daftar Isi", level=1)
    for i, item in enumerate([
        "Gambaran Umum Sistem",
        "Persiapan Operasional",
        "Alur Operasional Dasar",
        "Mode QC Sticker",
        "Mode Component Counter",
        "Troubleshooting & Penyesuaian Parameter",
        "Status Badge & Indikator",
    ], 1):
        add_para(doc, f"{i}. {item}", size=11)

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 1: GAMBARAN UMUM
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "1", "Gambaran Umum Sistem")

    add_para(doc, "QC Suite adalah sistem inspeksi visual berbasis AI untuk verifikasi komponen pada produk. Sistem menggunakan kamera dan model YOLO untuk mendeteksi, mengklasifikasi, dan menghitung komponen secara otomatis.", size=10.5)

    add_heading(doc, "1.1 Fitur Utama", level=2)
    add_bullet(doc, "Mode QC Sticker — Verifikasi sticker pada part menggunakan template matching dan OCR")
    add_bullet(doc, "Mode Component Counter — Penghitungan komponen berdasarkan ROI dan class target")
    add_bullet(doc, "Live View — Preview kamera real-time dengan overlay ROI dan deteksi")
    add_bullet(doc, "Counter Panel — Panel statistik accept/reject per sesi")
    add_bullet(doc, "PLC Integration — Kontrol clamp dan buzzer via Modbus RTU")
    add_bullet(doc, "Auto-reconnect — Kamera otomatis reconnect saat koneksi terputus")

    add_heading(doc, "1.2 Arsitektur Sistem", level=2)
    add_para(doc, "Sistem terdiri dari 3 komponen utama:", size=10.5)
    add_table(doc, ["Komponen", "Deskripsi"], [
        ["Backend (Flask)", "REST API, inference engine, session management, PLC worker"],
        ["Client Tkinter", "UI operator/admin, camera capture, frame upload"],
        ["Shared Contracts", "Template definitions, ROI geometry, dataclass contracts"],
    ])

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 2: PERSIAPAN OPERASIONAL
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "2", "Persiapan Operasional")

    add_heading(doc, "2.1 Persiapan Template (Admin)", level=2)
    add_para(doc, "Sebelum operator bisa menjalankan inspeksi, admin harus menyiapkan template melalui Admin Panel:", size=10.5)

    add_para(doc, "Langkah-langkah:", bold=True, size=10.5)
    add_bullet(doc, "Login sebagai Admin → Buka tab 'Presets'")
    add_bullet(doc, "Klik 'New Preset' → Isi form wizard")
    add_bullet(doc, "Pilih Validation Mode: 'QC Sticker' atau 'Component Counter'")
    add_bullet(doc, "Pilih Model YOLO yang sudah di-train")
    add_bullet(doc, "Atur Part Ready ROI — area untuk deteksi keberadaan part")
    add_bullet(doc, "Atur Sticker ROI — area untuk verifikasi sticker")
    add_bullet(doc, "Untuk Component Counter: tambah ROI + class targets (nama class + jumlah)")
    add_bullet(doc, "Capture Reference — ambil gambar referensi untuk gap detection")
    add_bullet(doc, "Capture Logo — ambil gambar logo untuk anti-reclamp (opsional)")
    add_bullet(doc, "Klik 'Save & Deploy Preset'")

    # Photo: Admin Panel
    add_photo_placeholder(doc, "Gambar 2.1 — Admin Panel, tab Presets")

    add_heading(doc, "2.2 Persiapan Sesi Operator", level=2)
    add_para(doc, "Setelah template di-deploy, operator perlu:", size=10.5)
    add_bullet(doc, "Login sebagai Operator")
    add_bullet(doc, "Pastikan kamera terhubung dan menampilkan gambar")
    add_bullet(doc, "Pilih template dari dropdown 'Template'")
    add_bullet(doc, "Pastikan status badge SERVER dan CAMERA hijau")
    add_bullet(doc, "Klik 'Start' untuk memulai inspeksi")

    # Photo: Operator Panel
    add_photo_placeholder(doc, "Gambar 2.2 — Operator Panel, tampilan awal")

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 3: ALUR OPERASIONAL DASAR
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "3", "Alur Operasional Dasar")

    add_heading(doc, "3.1 Start Inspeksi", level=2)
    add_para(doc, "Alur standar memulai inspeksi:", size=10.5)
    add_table(doc, ["Langkah", "Aksi", "Indikator"], [
        ["1", "Login sebagai Operator", "Context bar menampilkan nama operator"],
        ["2", "Pilih template dari dropdown", "Context bar menampilkan template aktif"],
        ["3", "Pastikan kamera aktif", "Badge CAMERA: CONNECTED (hijau)"],
        ["4", "Klik tombol 'Start'", "Decision banner berubah dari IDLE ke WAITING"],
        ["5", "Sistem otomatis start session", "Inference berjalan di background"],
    ])

    # Photo: Start Button
    add_photo_placeholder(doc, "Gambar 3.1 — Tombol Start di Operator Panel")

    add_heading(doc, "3.2 Proses Inspeksi", level=2)
    add_para(doc, "Setelah session aktif, sistem menjalankan loop berikut setiap frame:", size=10.5)
    add_table(doc, ["Tahap", "Deskripsi", "Output"], [
        ["Capture", "Ambil frame dari kamera", "Frame BGR"],
        ["Part Ready Check", "Cek apakah part ada di ROI", "part_ready: true/false"],
        ["Logo Check", "Cek apakah logo match (anti-reclamp)", "logo_skip: true/false"],
        ["Clamp", "PLC engage clamp saat part ready", "Clamp: ON/OFF"],
        ["Infer", "YOLO predict pada sticker ROI", "Deteksi + klasifikasi"],
        ["Validate", "Bandingkan hasil dengan target", "ACCEPT/REJECT"],
        ["Commit", "Simpan hasil ke database", "Counter +1"],
    ])

    # Photo: Live View
    add_photo_placeholder(doc, "Gambar 3.2 — Live View dengan overlay ROI dan deteksi")

    add_heading(doc, "3.3 Stop Inspeksi", level=2)
    add_para(doc, "Untuk menghentikan inspeksi:", size=10.5)
    add_bullet(doc, "Klik tombol 'Stop'")
    add_bullet(doc, "Sistem akan menghentikan inference dan camera")
    add_bullet(doc, "Session di-close di backend")
    add_bullet(doc, "Counter panel menampilkan statistik akhir")

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 4: MODE QC STICKER
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "4", "Mode QC Sticker")

    add_para(doc, "Mode QC Sticker digunakan untuk memverifikasi bahwa sticker yang terpasang pada part adalah sticker yang benar.", size=10.5)

    add_heading(doc, "4.1 Parameter Template", level=2)
    add_table(doc, ["Parameter", "Deskripsi", "Cara Atur"], [
        ["Expected Class", "Nama class sticker yang diharapkan", "Isi di wizard → Expected Class"],
        ["Sticker Code", "Kode OCR yang harus terbaca", "Isi di wizard → Sticker Code"],
        ["Use OCR", "Aktifkan verifikasi OCR", "Centang 'Use OCR verification'"],
        ["Min Confidence", "Threshold confidence minimum", "Isi di wizard → Confidence Threshold"],
        ["Max Tilt Degrees", "Batas kemiringan sticker (derajat)", "Isi di wizard → Max Tilt Degrees"],
        ["Gap Threshold", "Threshold match ratio part ready", "Isi di wizard → Gap Threshold"],
        ["Commit Stable Frames", "Jumlah frame stabil sebelum commit", "Default: 1"],
    ])

    add_heading(doc, "4.2 Alur Kerja", level=2)
    add_para(doc, "1. Part masuk ke area kamera → sistem mendeteksi via Part Ready ROI", size=10.5)
    add_para(doc, "2. Setelah part settled, PLC clamp engage → sticker terpasang", size=10.5)
    add_para(doc, "3. Sistem crop Sticker ROI → YOLO predict → deteksi class + confidence", size=10.5)
    add_para(doc, "4. Jika Use OCR aktif → baca teks sticker → bandingkan dengan Sticker Code", size=10.5)
    add_para(doc, "5. Hasil validasi → ACCEPT (hijau) atau REJECT (merah)", size=10.5)
    add_para(doc, "6. Setelah ACCEPT → blackout timer → operator bisa ganti part", size=10.5)

    # Photo: QC Sticker Flow
    add_photo_placeholder(doc, "Gambar 4.1 — Alur inspeksi QC Sticker (Part Ready → Clamp → Infer → Validate)")

    add_heading(doc, "4.3 Reject Reasons", level=2)
    add_table(doc, ["Kode", "Penyebab", "Solusi"], [
        ["WRONG_TYPE", "Class sticker tidak sesuai expected", "Ganti sticker yang benar"],
        ["WRONG_TEXT", "OCR text tidak sesuai sticker code", "Periksa posisi sticker"],
        ["LOW_CLASS_CONF", "Confidence class di bawah threshold", "Perbaiki pencahayaan / posisi kamera"],
        ["LOW_ROI_CONF", "Confidence ROI terlalu rendah", "Periksa ROI sticker"],
        ["OUT_OF_ANGLE", "Kemiringan sticker melebihi batas", "Pasang sticker lebih rata"],
        ["PART_NOT_READY", "Part tidak terdeteksi di ROI", "Pastikan part masuk ROI"],
    ])

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 5: MODE COMPONENT COUNTER
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "5", "Mode Component Counter")

    add_para(doc, "Mode Component Counter digunakan untuk memverifikasi jumlah dan jenis komponen yang terpasang pada produk. Sistem menggunakan multiple ROI dan YOLO untuk menghitung komponen per area.", size=10.5)

    add_heading(doc, "5.1 Konfigurasi ROI", level=2)
    add_para(doc, "Setiap ROI memiliki:", size=10.5)
    add_bullet(doc, "Name — Nama ROI (contoh: 'ROI A', 'Bolt Area')")
    add_bullet(doc, "Geometry — Posisi (x, y) dan ukuran (w, h) dalam koordinat fractional 0-1")
    add_bullet(doc, "Classes — Daftar class target dengan jumlah yang diharapkan")
    add_bullet(doc, "Strict Foreign Class — Jika aktif, class terdeteksi selain target akan menyebabkan REJECT")

    # Photo: Component ROI Config
    add_photo_placeholder(doc, "Gambar 5.1 — Konfigurasi Component ROI di Admin Panel")

    add_heading(doc, "5.2 Alur Kerja", level=2)
    add_para(doc, "1. Part masuk ke area kamera → sistem mendeteksi via Part Ready ROI", size=10.5)
    add_para(doc, "2. Setelah part settled, PLC clamp engage", size=10.5)
    add_para(doc, "3. Sistem crop setiap Component ROI → susun jadi montage grid", size=10.5)
    add_para(doc, "4. 1× YOLO predict pada montage → deteksi semua komponen", size=10.5)
    add_para(doc, "5. Remap detections ke koordinat frame penuh → assign ke ROI", size=10.5)
    add_para(doc, "6. Voting K-frame → ambil modus count per class", size=10.5)
    add_para(doc, "7. Evaluasi: semua class exact + total match → ACCEPT, selain itu → REJECT", size=10.5)

    # Photo: Montage Flow
    add_photo_placeholder(doc, "Gambar 5.2 — Proses montage: crop ROI → grid → 1× infer → remap")

    add_heading(doc, "5.3 Reject Reasons", level=2)
    add_table(doc, ["Kode", "Penyebab", "Solusi"], [
        ["COMPONENT_COUNT_MISMATCH", "Jumlah komponen tidak sesuai target", "Periksa komponen yang kurang/berlebih"],
        ["UNEXPECTED_COMPONENT", "Ada komponen asing atau jumlah total tidak cocok", "Periksa apakah ada komponen yang tidak seharusnya ada"],
        ["NO_COMPONENT_ROIS", "Tidak ada ROI yang didefinisikan di template", "Admin perlu menambah ROI di preset"],
    ])

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 6: TROUBLESHOOTING
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "6", "Troubleshooting & Penyesuaian Parameter")

    add_heading(doc, "6.1 Match Ratio Turun", level=2)
    add_para(doc, "Jika gap match ratio turun di bawah threshold (default 0.85):", size=10.5)
    add_table(doc, ["Penyebab", "Solusi"], [
        ["Pencahayaan berubah", "Atur exposure kamera atau tambah lampu"],
        ["Kamera bergeser/getar", "Kencangkan mount kamera, aktifkan logo anti-reclamp"],
        ["Reference patch outdated", "Capture ulang reference di Admin Panel"],
        ["Part berbeda", "Pastikan part yang sama digunakan untuk reference"],
    ])

    add_heading(doc, "6.2 Class Sticker Berubah", level=2)
    add_para(doc, "Jika class sticker berubah (model update atau produk baru):", size=10.5)
    add_bullet(doc, "Buka Admin Panel → Tab Presets")
    add_bullet(doc, "Edit preset yang relevan")
    add_bullet(doc, "Update 'Expected Class' dengan class baru")
    add_bullet(doc, "Jika menggunakan OCR, update 'Sticker Code'")
    add_bullet(doc, "Klik 'Save & Deploy Preset'")
    add_bullet(doc, "Operator perlu refresh template list (klik 'Refresh')")

    add_heading(doc, "6.3 Parameter yang Sering Diubah", level=2)
    add_table(doc, ["Parameter", "Lokasi", "Efek"], [
        ["Gap Threshold", "Admin → Preset Wizard → Gap Threshold", "Semakin tinggi = semakin ketat deteksi part"],
        ["Confidence Threshold", "Admin → Preset Wizard → Confidence Threshold", "Filter deteksi YOLO dengan confidence minimum"],
        ["Max Tilt Degrees", "Admin → Preset Wizard → Max Tilt Degrees", "Batas kemiringan sticker yang diterima"],
        ["Commit Stable Frames", "Admin → Sticker Config", "Jumlah frame stabil sebelum commit (default 1)"],
        ["Phase Next Part Delay", "Environment variable (ms)", "Waktu tunggu setelah ACCEPT sebelum part baru"],
        ["Inference Interval", "Environment variable (ms)", "Interval inference (0 = setiap frame)"],
    ])

    add_heading(doc, "6.4 Kamera Terputus", level=2)
    add_para(doc, "Jika kamera terputus:", size=10.5)
    add_bullet(doc, "Sistem otomatis menampilkan status 'CAMERA: RECONNECTING' (kuning)")
    add_bullet(doc, "Auto-reconnect berjalan — sistem mencoba reconnect setiap 2 detik")
    add_bullet(doc, "Inference di-skip selama reconnecting")
    add_bullet(doc, "Setelah reconnect, inference otomatis resume")
    add_bullet(doc, "Jika gagal setelah 5 kali coba, badge berubah 'CAMERA: DISCONNECTED' (merah)")

    add_heading(doc, "6.5 Session Error", level=2)
    add_para(doc, "Jika session error:", size=10.5)
    add_bullet(doc, "Cek badge SERVER — jika merah, backend tidak tersedia")
    add_bullet(doc, "Cek badge SESSION — jika merah, session tidak aktif")
    add_bullet(doc, "Klik 'Stop' lalu 'Start' untuk restart session")
    add_bullet(doc, "Jika masih error, cek log di backend console")

    add_page_break(doc)

    # ═══════════════════════════════════════════════════════════
    # BAB 7: STATUS BADGE & INDIKATOR
    # ═══════════════════════════════════════════════════════════
    add_section_header(doc, "7", "Status Badge & Indikator")

    add_heading(doc, "7.1 System Status Badges", level=2)
    add_para(doc, "Badge di status strip menunjukkan kondisi sistem:", size=10.5)
    add_table(doc, ["Badge", "Hijau", "Merah", "Kuning"], [
        ["SERVER", "Backend connected", "Backend unreachable", "—"],
        ["CAMERA", "Camera active", "Camera disconnected", "Reconnecting"],
        ["SESSION", "Session active", "No session", "—"],
        ["DB", "Database OK", "DB error", "—"],
        ["EVENT", "Event running", "No event", "—"],
        ["PLC", "PLC connected", "PLC error", "—"],
    ])

    # Photo: Status Badges
    add_photo_placeholder(doc, "Gambar 7.1 — Status badges di Operator Panel")

    add_heading(doc, "7.2 Decision Banner", level=2)
    add_para(doc, "Banner utama di tengah layar:", size=10.5)
    add_bullet(doc, "IDLE (abu-abu) — Sistem menunggu, belum ada session aktif")
    add_bullet(doc, "WAITING (biru) — Session aktif, menunggu event inspeksi pertama")
    add_bullet(doc, "ACCEPT (hijau) — Inspeksi berhasil, komponen sesuai")
    add_bullet(doc, "REJECT (merah) — Inspeksi gagal, ada ketidaksesuaian")

    add_heading(doc, "7.3 Counter Panel", level=2)
    add_para(doc, "Panel counter di sisi kanan menampilkan:", size=10.5)
    add_bullet(doc, "Total — Total frame yang diproses")
    add_bullet(doc, "Accept — Jumlah inspeksi berhasil")
    add_bullet(doc, "Reject — Jumlah inspeksi gagal")
    add_bullet(doc, "Reject Breakdown — Detail reject per kode (WRONG_TYPE, LOW_CLASS_CONF, dll)")

    # Photo: Counter Panel
    add_photo_placeholder(doc, "Gambar 7.2 — Counter Panel dan Decision Banner")

    add_heading(doc, "7.4 Live View Overlay", level=2)
    add_para(doc, "Overlay pada live view:", size=10.5)
    add_bullet(doc, "Part Ready ROI (oranye) — Area deteksi part")
    add_bullet(doc, "Sticker ROI (kuning) — Area verifikasi sticker")
    add_bullet(doc, "Component ROI (hijau) — Area penghitungan komponen")
    add_bullet(doc, "Crosshair (putih) — Titik center expected")
    add_bullet(doc, "Bounding box deteksi — Kotak hasil YOLO predict")

    add_page_break(doc)

    # ── Footer ──
    add_para(doc, "", size=10)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("QC Suite Operator Manual Book — PT ASKI — 2026")
    run.font.size = Pt(9)
    run.font.color.rgb = GRAY
    run.font.name = 'Calibri'

    os.makedirs(GUIDEBOOK_DIR, exist_ok=True)
    docx_path = os.path.join(GUIDEBOOK_DIR, "QC_Suite_Operator_Manual.docx")
    doc.save(docx_path)
    print(f"DOCX saved: {docx_path}")
    return docx_path


# ═══════════════════════════════════════════════════════════════════
# PPTX GENERATION
# ═══════════════════════════════════════════════════════════════════

def add_pptx_title_slide(prs, title, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(10), PptxInches(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = P_BLUE_DARK
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(2), PptxInches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "[ LOGO ASKI ]"
    p.font.size = PptxPt(12)
    p.font.color.rgb = P_WHITE
    p.font.name = 'Calibri'
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = P_WHITE
    p.font.name = 'Calibri'
    p = tf.add_paragraph()
    p.text = subtitle_text
    p.font.size = PptxPt(20)
    p.font.color.rgb = P_BLUE_LIGHT
    p.font.name = 'Calibri'
    shape2 = slide.shapes.add_shape(1, PptxInches(0), PptxInches(5), PptxInches(10), PptxInches(2.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = P_BLUE_MID
    shape2.line.fill.background()
    txBox2 = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(5.5), PptxInches(9), PptxInches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "PT ASKI  —  Versi 1.0  —  2026"
    p2.font.size = PptxPt(18)
    p2.font.bold = True
    p2.font.color.rgb = P_WHITE
    p2.font.name = 'Calibri'


def add_pptx_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shape = slide.shapes.add_shape(1, PptxInches(0.3), PptxInches(0.2), PptxInches(9.4), PptxInches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = P_BLUE_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(24)
    p.font.bold = True
    p.font.color.rgb = P_WHITE
    p.font.name = 'Calibri'
    txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.2), PptxInches(9), PptxInches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = PptxPt(14)
        p.font.color.rgb = P_DARK
        p.font.name = 'Calibri'
        p.space_after = PptxPt(6)


def add_pptx_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shape = slide.shapes.add_shape(1, PptxInches(0.3), PptxInches(0.2), PptxInches(9.4), PptxInches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = P_BLUE_DARK
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptxPt(24)
    p.font.bold = True
    p.font.color.rgb = P_WHITE
    p.font.name = 'Calibri'
    rows_count = len(rows) + 1
    cols_count = len(headers)
    tbl_shape = slide.shapes.add_table(rows_count, cols_count, PptxInches(0.5), PptxInches(1.2), PptxInches(9), PptxInches(4.5))
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = PptxPt(12)
            p.font.bold = True
            p.font.color.rgb = P_WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = P_BLUE_MID
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = PptxPt(10)
                p.font.color.rgb = P_DARK
                p.font.name = 'Calibri'
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = P_BLUE_BG


def generate_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)

    # ── Helper: add step-by-step slide ──
    def add_step_slide(title, steps, note=None):
        """Add a compact step-by-step slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Title bar
        shape = slide.shapes.add_shape(1, PptxInches(0.3), PptxInches(0.15), PptxInches(9.4), PptxInches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = P_BLUE_DARK
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptxPt(22)
        p.font.bold = True
        p.font.color.rgb = P_WHITE
        p.font.name = 'Calibri'
        # Steps
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.0), PptxInches(9), PptxInches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, step in enumerate(steps):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # Step number (bold, blue)
            run_num = p.add_run()
            run_num.text = "  %d. " % (i + 1)
            run_num.font.bold = True
            run_num.font.size = PptxPt(13)
            run_num.font.color.rgb = P_BLUE_MID
            run_num.font.name = 'Calibri'
            # Step description
            run_desc = p.add_run()
            run_desc.text = step
            run_desc.font.size = PptxPt(13)
            run_desc.font.color.rgb = P_DARK
            run_desc.font.name = 'Calibri'
            p.space_after = PptxPt(4)
        # Note at bottom
        if note:
            p_note = tf.add_paragraph()
            p_run = p_note.add_run()
            p_run.text = "📷 %s" % note
            p_run.font.size = PptxPt(11)
            p_run.font.color.rgb = P_GRAY
            p_run.font.name = 'Calibri'
            p_run.font.italic = True

    # ── Helper: add compact table slide ──
    def add_tbl(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_shape(1, PptxInches(0.3), PptxInches(0.15), PptxInches(9.4), PptxInches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = P_BLUE_DARK
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptxPt(22)
        p.font.bold = True
        p.font.color.rgb = P_WHITE
        p.font.name = 'Calibri'
        tbl_shape = slide.shapes.add_table(len(rows)+1, len(headers), PptxInches(0.5), PptxInches(1.0), PptxInches(9), PptxInches(4.5))
        tbl = tbl_shape.table
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = PptxPt(11)
                p.font.bold = True
                p.font.color.rgb = P_WHITE
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = P_BLUE_MID
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri+1, ci)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = PptxPt(10)
                    p.font.color.rgb = P_DARK
                    p.font.name = 'Calibri'
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = P_BLUE_BG

    # ═══════════════════════════════════════════════════════════
    # SLIDES
    # ═══════════════════════════════════════════════════════════

    # Slide 1: Cover
    add_pptx_title_slide(prs, "QC SUITE", "Operator Manual Book — Sistem Inspeksi QC Sticker & Component Counter")

    # Slide 2: TOC
    add_pptx_content_slide(prs, "Daftar Isi", [
        "1. Gambaran Umum",
        "2. Persiapan Operasional",
        "3. Alur Operasional Dasar",
        "4. Mode QC Sticker",
        "5. Mode Component Counter",
        "6. Troubleshooting",
        "7. Status Badge & Indikator",
    ])

    # Slide 3: Overview
    add_pptx_content_slide(prs, "1. Gambaran Umum", [
        "QC Suite = sistem inspeksi visual berbasis AI",
        "Kamera + YOLO untuk deteksi & klasifikasi",
        "Dua mode: QC Sticker & Component Counter",
        "Arsitektur: Backend (Flask) + Client (Tkinter)",
        "Fitur: Live View, Counter Panel, PLC, Auto-reconnect",
    ])

    # Slide 4: Persiapan Admin
    add_step_slide("2.1 Persiapan Template (Admin)", [
        "Login Admin → Tab 'Presets' → 'New Preset'",
        "Pilih Validation Mode: QC Sticker / Component Counter",
        "Pilih Model YOLO → Atur Part Ready ROI & Sticker ROI",
        "Untuk Component Counter: tambah ROI + class targets",
        "Capture Reference → Capture Logo (opsional)",
        "Klik 'Save & Deploy Preset'",
    ], "Admin Panel → tab Presets")

    # Slide 5: Persiapan Operator
    add_step_slide("2.2 Persiapan Sesi Operator", [
        "Login sebagai Operator",
        "Pastikan kamera terhubung (badge CAMERA: hijau)",
        "Pilih template dari dropdown 'Template'",
        "Pastikan badge SERVER & CAMERA hijau",
        "Klik 'Start' untuk memulai inspeksi",
    ], "Operator Panel → tampilan awal")

    # Slide 6: Start
    add_tbl("3.1 Start Inspeksi",
        ["#", "Aksi", "Indikator"],
        [
            ["1", "Login sebagai Operator", "Context bar: nama operator"],
            ["2", "Pilih template", "Context bar: template aktif"],
            ["3", "Pastikan kamera aktif", "Badge CAMERA: CONNECTED"],
            ["4", "Klik 'Start'", "Decision banner: WAITING"],
            ["5", "Session aktif", "Inference berjalan di background"],
        ])

    # Slide 7: Proses per Frame
    add_tbl("3.2 Proses Inspeksi per Frame",
        ["Tahap", "Deskripsi", "Output"],
        [
            ["Capture", "Ambil frame dari kamera", "Frame BGR"],
            ["Part Ready", "Cek part di Part Ready ROI", "part_ready: true/false"],
            ["Logo Check", "Cek logo match (anti-reclamp)", "logo_skip: true/false"],
            ["Clamp", "PLC engage clamp", "Clamp: ON/OFF"],
            ["Infer", "YOLO predict pada ROI", "Deteksi + klasifikasi"],
            ["Validate", "Bandingkan dengan target", "ACCEPT / REJECT"],
            ["Commit", "Simpan ke database", "Counter +1"],
        ])

    # Slide 8: QC Sticker
    add_step_slide("4. Mode QC Sticker", [
        "Verifikasi sticker pada part menggunakan template matching + OCR",
        "Parameter: Expected Class, Sticker Code, Confidence Threshold",
        "Alur: Part masuk → settle → clamp → infer → validate → commit",
        "Setelah ACCEPT: blackout timer → operator ganti part",
        "Reject jika: WRONG_TYPE, WRONG_TEXT, LOW_CONF, OUT_OF_ANGLE",
    ], "Alur inspeksi QC Sticker")

    # Slide 9: QC Params
    add_tbl("4.1 Parameter Template QC Sticker",
        ["Parameter", "Deskripsi"],
        [
            ["Expected Class", "Nama class sticker yang diharapkan"],
            ["Sticker Code", "Kode OCR yang harus terbaca"],
            ["Use OCR", "Aktifkan verifikasi OCR"],
            ["Min Confidence", "Threshold confidence minimum"],
            ["Max Tilt Degrees", "Batas kemiringan sticker"],
            ["Gap Threshold", "Threshold match ratio part ready"],
        ])

    # Slide 10: Component Counter
    add_step_slide("5. Mode Component Counter", [
        "Penghitungan komponen berdasarkan ROI dan class target",
        "Setiap ROI: Name, Geometry, Classes (class_name + count)",
        "Alur: crop ROI → montage grid → 1× YOLO → remap → voting",
        "ACCEPT jika semua class exact + total match",
        "Reject jika: COMPONENT_COUNT_MISMATCH, UNEXPECTED_COMPONENT",
    ], "Konfigurasi Component ROI di Admin Panel")

    # Slide 11: Montage
    add_step_slide("5.1 Proses Montage", [
        "1. Crop setiap Component ROI dari frame penuh",
        "2. Susun crops jadi grid (montage)",
        "3. 1× YOLO predict pada montage",
        "4. Remap detections ke koordinat frame penuh",
        "5. Assign ke ROI berdasarkan center point",
        "6. Voting K-frame → modus count per class",
    ], "Proses montage → grid → infer → remap")

    # Slide 12: Troubleshooting
    add_tbl("6.1 Troubleshooting",
        ["Masalah", "Penyebab", "Solusi"],
        [
            ["Match ratio turun", "Pencahayaan berubah / kamera geser", "Atur exposure / capture ulang reference"],
            ["Class sticker berubah", "Model update / produk baru", "Update Expected Class di Admin Panel"],
            ["Kamera terputus", "Kabel longgar / driver error", "Cek kabel / sistem auto-reconnect"],
            ["Session error", "Backend down / network issue", "Stop → Start ulang / cek log"],
        ])

    # Slide 13: Parameter
    add_tbl("6.2 Parameter yang Sering Diubah",
        ["Parameter", "Lokasi", "Efek"],
        [
            ["Gap Threshold", "Admin → Preset Wizard", "Semakin tinggi = semakin ketat"],
            ["Confidence Threshold", "Admin → Preset Wizard", "Filter deteksi YOLO"],
            ["Max Tilt Degrees", "Admin → Preset Wizard", "Batas kemiringan sticker"],
            ["Commit Stable Frames", "Admin → Sticker Config", "Frame stabil sebelum commit"],
        ])

    # Slide 14: Status Badges
    add_tbl("7.1 System Status Badges",
        ["Badge", "Hijau", "Merah", "Kuning"],
        [
            ["SERVER", "Connected", "Unreachable", "—"],
            ["CAMERA", "Active", "Disconnected", "Reconnecting"],
            ["SESSION", "Active", "No session", "—"],
            ["DB", "OK", "Error", "—"],
            ["EVENT", "Running", "No event", "—"],
            ["PLC", "Connected", "Error", "—"],
        ])

    # Slide 15: Banner & Counter
    add_step_slide("7.2 Decision Banner & Counter", [
        "IDLE (abu-abu) — Sistem menunggu",
        "WAITING (biru) — Session aktif, menunggu event",
        "ACCEPT (hijau) — Inspeksi berhasil",
        "REJECT (merah) — Inspeksi gagal",
        "Counter Panel: Total / Accept / Reject / Breakdown",
    ], "Counter Panel dan Decision Banner")

    # Slide 16: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(10), PptxInches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = P_BLUE_DARK
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.5), PptxInches(8), PptxInches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Terima Kasih"
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = P_WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "QC Suite Operator Manual Book"
    p.font.size = PptxPt(20)
    p.font.color.rgb = P_BLUE_LIGHT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "PT ASKI — 2026"
    p.font.size = PptxPt(14)
    p.font.color.rgb = P_WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    pptx_path = os.path.join(GUIDEBOOK_DIR, "QC_Suite_Operator_Manual.pptx")
    prs.save(pptx_path)
    print(f"PPTX saved: {pptx_path}")
    return pptx_path


if __name__ == "__main__":
    docx_path = generate_docx()
    pptx_path = generate_pptx()
    print(f"\nGuidebook generated:")
    print(f"  DOCX: {docx_path}")
    print(f"  PPTX: {pptx_path}")
